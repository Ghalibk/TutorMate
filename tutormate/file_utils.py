import PyPDF2
import docx
from pptx import Presentation

def extract_text_from_file(file_path):
    """Determine the file type and extract text."""
    try:
        if file_path.endswith('.pdf'):
            return extract_text_from_pdf(file_path)
        elif file_path.endswith('.docx'):
            return extract_text_from_docx(file_path)
        elif file_path.endswith('.pptx'):
            return extract_text_from_pptx(file_path)
        else:
            raise ValueError("Unsupported file format")
    except Exception as e:
        raise ValueError(f"Error extracting text from file: {e}")

def extract_text_from_pdf(file_path):
    """Extract text from a PDF file."""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text()
    except Exception as e:
        raise ValueError(f"Error extracting text from PDF: {e}")
    return text

def extract_text_from_docx(file_path):
    """Extract text from a Word document."""
    text = ""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        raise ValueError(f"Error extracting text from DOCX: {e}")
    return text

def extract_text_from_pptx(file_path):
    """Extract text from a PowerPoint file."""
    text = ""
    try:
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"
    except Exception as e:
        raise ValueError(f"Error extracting text from PPTX: {e}")
    return text
