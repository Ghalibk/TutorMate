'''
    Views to access Microsoft Graph API for user details.
'''

# Django Imports
from django.conf import settings
from django.core.files.base import ContentFile
from django.core.files.storage import default_storage
from django.http import HttpResponse, JsonResponse,  FileResponse
from django.shortcuts import render, redirect
from django.views.decorators.csrf import csrf_exempt

# REST Framework Imports
from rest_framework import status
from rest_framework import serializers
from rest_framework.decorators import api_view
from rest_framework.response import Response
from rest_framework.views import APIView

# Application Imports
from .models import Course, Enroll, User, Todo, Module, Quiz, Flashcard, Bulletpoint, Fullsummary
from .sync import userDataSync
from .utils import generate_quiz, generate_flashcard, generate_bulletpoints, generate_full_summary, generate_todo_help

# External Imports
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from datetime import datetime
import json
import os
import requests
import threading
from canvas import *

def get_graph_token():
    '''Get token from Microsoft AAD URL.'''
    try:
        url = settings.AD_URL

        headers = {'Content-Type': 'application/x-www-form-urlencoded', 'Accept': 'application/json'}

        data = {
            'grant_type': 'client_credentials',
            'client_id': settings.CLIENT_ID,
            'client_secret': settings.CLIENT_SECRET,
            'scope': 'https://graph.microsoft.com/.default',
        }

        return requests.post(url=url, headers=headers, data=data).json()
    except:
        return None

def login_successful(request):
    if request.user.is_authenticated:
        ''' Get User details from Microsoft Graph APIs.'''
        graph_token = get_graph_token()
        print(graph_token)
        
        if graph_token:
            # Step 1: Fetch User Details
            user_url = f"https://graph.microsoft.com/v1.0/users/{request.user.username}"
            headers = {
                "Authorization": f"Bearer {graph_token['access_token']}",
                "Content-Type": "application/json",
            }
            user_response = requests.get(url=user_url, headers=headers).json()
            print("User Details:", user_response)

            # Step 2: Fetch User Profile Picture
            photo_url = f"https://graph.microsoft.com/v1.0/users/{request.user.username}/photo/$value"
            photo_response = requests.get(url=photo_url, headers=headers)

            if photo_response.status_code == 200:
                # Profile picture retrieved successfully
                with open("profile_picture.jpg", "wb") as f:
                    f.write(photo_response.content)
                print("Profile picture saved as 'profile_picture.jpg'")
                return HttpResponse("User profile and photo retrieved successfully.")
            else:
                # Handle errors (e.g., no profile picture exists)
                print(f"Failed to retrieve profile picture: {photo_response.status_code}")
                return HttpResponse("Failed to retrieve profile picture.", status=photo_response.status_code)
        else:
            return HttpResponse("Failed to retrieve token.", status=400)
    else:
        return HttpResponse("User not authenticated.", status=401)
    
def get_user_photo(request):
    if request.user.is_authenticated:
        try:
            user = User.objects.get(user_id=request.session["id"])
            if user.profile_picture:
                print(f"Serving profile picture: {user.profile_picture.path}")
                return FileResponse(user.profile_picture, content_type="image/jpeg")
            else:
                print("No profile picture found.")
                return HttpResponse("No profile picture available.", status=404)
        except User.DoesNotExist:
            return HttpResponse("User not found.", status=404)
    print("User not authenticated.")
    return HttpResponse("User not authenticated.", status=401)


def get_user_info(request):
    id = ""
    name = ""
    email = ""
    if request.user.is_authenticated:
        id = request.session["id"]
        name = request.session["first_name"] + ' ' + request.session["last_name"]
        email = request.session["email"]
    data = {
        'id': id,
        'name': name,
        'email': email,
    }
    return JsonResponse(data)


def login(request):
    # Redirect to the ADFS login page
    return redirect(settings.LOGIN_URL) 

@api_view(['GET'])
def fetch_courses(request):
    # Ensure user is authenticated
    if not request.user.is_authenticated:
        return Response({"status": "error", "message": "User is not authenticated"}, status=401)

    try:
        # Get user info from the session
        user = User.objects.get(user_id=request.session["id"])

        # Fetch courses from the database
        courses = Course.objects.all()  # You can add any filters if necessary
        course_data = []

        for course in courses:
            # Fetch the enrollment data for the user and course using user_id and course_id
            enrollment = Enroll.objects.filter(user=user, course=course).first()

            if enrollment:
                course_details = {
                    "id": course.course_id,
                    "name": course.name,
                    "term_name": course.term_name,
                    "image_url": course.image_url,
                    "letter_grade": enrollment.letter_grade,
                    "overall_grade": enrollment.percent,
                }
                course_data.append(course_details)

        if course_data:
            print(course_data)
            return Response(course_data)
        else:
            return Response({"status": "error", "message": "No courses found."}, status=404)

    except:
        return Response({"status": "error", "message": "User not found in the database."}, status=404)

def check_canvas_token(request):
    if request.user.is_authenticated:
        try:
            # Fetch the User object for the logged-in user
            user = User.objects.get(user_id=request.session["id"])
            # Check if canvas_token is null or initialized
            token_status = "initialized" if user.canvas_token else "null"
            if(token_status == "initialized"):
                if(not CanvasConnexion(user.canvas_token).is_token_valid()):
                    token_status = "not valid"
            response = {
                "status": "success",
                "token_status": token_status,
            }
        except User.DoesNotExist:
            # Handle case where the User does not exist in the database
            response = {
                "status": "error",
                "message": "User not found in the database."
            }

    return JsonResponse(response)

def validate_canvas_token(request):
    if request.method == "POST":
        try:
            # Parse the token from the request body
            data = json.loads(request.body)
            canvas_token = data.get("token")

            if not canvas_token:
                return JsonResponse({"status": "error", "message": "Token is required."}, status=400)

            # Validate the token using CanvasConnexion
            is_valid = CanvasConnexion(canvas_token).is_token_valid()

            if is_valid:
                # Update the user's canvas_token in the database
                if request.user.is_authenticated:
                    try:
                        user = User.objects.get(user_id=request.session["id"])
                        user.canvas_token = canvas_token
                        user.save(update_fields=["canvas_token"])
                        thread = threading.Thread(target=userDataSync, args=(request.session["id"],))
                        thread.start()
                        graph_token = get_graph_token()
                        if graph_token:
                            photo_url = f"https://graph.microsoft.com/v1.0/users/{request.user.username}/photo/$value"
                            headers = {"Authorization": f"Bearer {graph_token['access_token']}"}
                            photo_response = requests.get(url=photo_url, headers=headers)
                            
                            if photo_response.status_code == 200:
                                try:
                                    # Ensure content is bytes
                                    content = photo_response.content
                                    if isinstance(content, memoryview):
                                        content = content.tobytes()

                                    # Save the profile picture properly
                                    user.profile_picture.save(
                                        f"{request.session['id']}.jpg",  # File name
                                        ContentFile(content),  # Content as a file
                                        save=True  # Automatically save the user instance
                                    )

                                    # Debugging
                                    print(f"Profile picture saved: {user.profile_picture}")  # Should print the file path, not <memory>
                                except Exception as e:
                                    print(f"Error saving profile picture: {e}")
                            else:
                                print(f"{request.session['id']} - Failed to fetch photo: {photo_response.status_code}")
                        else:
                            print(f"Graph token is missing for user {request.session['id']}")
                        return JsonResponse({"status": "success", "message": "Token is valid."})
                    except User.DoesNotExist:
                        return JsonResponse({"status": "error", "message": "User not found in the database."}, status=404)
                else:
                    return JsonResponse({"status": "error", "message": "User is not authenticated."}, status=401)
            else:
                return JsonResponse({"status": "error", "message": "Token is not valid."}, status=400)

        except json.JSONDecodeError:
            return JsonResponse({"status": "error", "message": "Invalid JSON format."}, status=400)

    return JsonResponse({"status": "error", "message": "Invalid request method."}, status=405)

def react_view(request, path=None):
    # Serve the React app's index.html for all frontend routes
    index_path = os.path.join(settings.BASE_DIR, 'frontend/tutormate/dist', 'index.html')
    with open(index_path, 'r') as f:
        return HttpResponse(f.read(), content_type='text/html')
    
@api_view(['GET'])
def fetch_todo(request):
    print("Fetching To-Do tasks...")

    if not request.user.is_authenticated:
        return Response({"status": "error", "message": "User is not authenticated"}, status=401)

    try:
        # Get the user ID from the session
        user_id = request.session.get("id")
        if not user_id:
            return Response({"status": "error", "message": "User ID is missing in session"}, status=400)

        # Fetch the courses the user is enrolled in from the Enroll table
        enrolled_courses = Enroll.objects.filter(user_id=user_id).values_list('course__course_id', flat=True)
        if not enrolled_courses:
            return Response({"status": "success", "message": "No courses found for the user", "todos": []}, status=200)

        # Fetch To-Dos for the courses the user is enrolled in
        todo_items = Todo.objects.filter(course__course_id__in=enrolled_courses)
        print("Filtered tasks:", todo_items)

        # List to store the formatted tasks
        formatted_items = []

        # Loop through each todo item and read the corresponding file
        for item in todo_items:
            file_path = f"./todo/{item.course.course_id}/{item.todo_id}.txt"

            # Check if the file exists
            if os.path.exists(file_path):
                with open(file_path, 'r') as file:
                    description = file.read()  # Read the content of the file
            else:
                description = "No description available"  # Default message if file does not exist

            # Check if 'due_date' is a string and convert it to datetime if necessary
            if isinstance(item.due_date, str):
                try:
                    due_date = datetime.strptime(item.due_date, "%Y-%m-%dT%H:%M:%SZ")  # Adjust format as needed
                except ValueError:
                    due_date = None
            else:
                due_date = item.due_date  # Assume it's already a datetime object

            # Create a formatted item with the description from the file
            formatted_items.append({
                "assignment_name": item.assignment_name,
                "due_date": due_date.strftime("%Y-%m-%dT%H:%M:%S") if due_date else "No Due Date",
                "assignment_url": item.assignment_url,
                "description": description,  # Adding the file content as description
            })

        # Return the formatted tasks as JSON
        return Response(formatted_items, status=200)

    except Exception as e:
        print(f"Error: {e}")
        return Response({"status": "error", "message": str(e)}, status=500)
    
def get_course_name(request):
    course_id = request.GET.get("courseid")
    if not course_id:
        return JsonResponse({"error": "Course ID is required"}, status=400)

    try:
        # Query the database for the course
        course = Course.objects.get(course_id=course_id)
        return JsonResponse({"name": course.name})
    except Course.DoesNotExist:
        return JsonResponse({"error": "Course not found"}, status=404)
    
def get_modules(request):
    course_id = request.GET.get("courseid")
    if not course_id:
        return JsonResponse({"error": "Course ID is required"}, status=400)

    # Query the database for modules linked to the specified course
    modules = Module.objects.filter(course_id=course_id).values("id", "file_name")
    
    if not modules.exists():
        return JsonResponse({"modules": []}, status=404)  # No modules found

    return JsonResponse({"modules": list(modules)}, status=200)

def get_todo(request):
    course_id = request.GET.get("courseid")
    if not course_id:
        return JsonResponse({"status": "error", "message": "Course ID is required"}, status=400)

    try:
        todos = Todo.objects.filter(course__course_id=course_id)

        if not todos.exists():
            # Return a response indicating no To-Do items found
            return JsonResponse(
                {"status": "success", "message": "No To-Do items found", "todos": []},
                status=200,
            )

        # Serialize todo_id and assignment_name
        todo_list = [{"todo_id": todo.todo_id, "assignment_name": todo.assignment_name} for todo in todos]

        return JsonResponse(
            {"status": "success", "todos": todo_list},
            status=200,
        )

    except Exception as e:
        print(f"Error fetching To-Do items: {e}")
        return JsonResponse(
            {"status": "error", "message": "An error occurred while fetching To-Do items"},
            status=500,
        )


def extract_text_from_file(file_path):
    """Extract text from PDF, PPTX, or DOCX."""
    try:
        _, file_extension = os.path.splitext(file_path)

        if file_extension.lower() == ".pdf":
            text = ""
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text()
            return text

        elif file_extension.lower() == ".pptx":
            text = ""
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text += paragraph.text
            return text

        elif file_extension.lower() == ".docx":
            text = ""
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text
            return text

        else:
            raise ValueError("Unsupported file type")

    except Exception as e:
        print(f"Error extracting text: {e}")
        return None

def save_json_to_file(directory, filename, data):
    # Ensure the directory exists
    os.makedirs(directory, exist_ok=True)
    
    # Write the JSON data to the specified file
    with open(os.path.join(directory, filename), "w") as file:
        json.dump(data, file, indent=4)

@csrf_exempt
def generate_quiz_view(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body)
            course_id = data.get("course_id")
            module_name = data.get("module_name")
            difficulty = data.get("difficulty")
            num_questions = data.get("num_questions")

            # Validate inputs
            if not course_id or not module_name or not difficulty or not num_questions:
                return JsonResponse({"status": "error", "message": "All fields are required."}, status=400)
            
            if difficulty not in ["easy", "medium", "difficult"]:
                return JsonResponse({"status": "error", "message": "Invalid difficulty level."}, status=400)

            num_questions = int(num_questions)
            if num_questions < 1 or num_questions > 50:
                return JsonResponse({"status": "error", "message": "Number of questions must be between 1 and 50."}, status=400)

            extracted_text = extract_text_from_file(f"./modules/{course_id}/{module_name}")
            if not extracted_text:
                return JsonResponse({"status": "error", "message": "Failed to extract text from the file."}, status=400)

            quiz_json_string = generate_quiz(difficulty, num_questions, extracted_text)
            quiz_data = json.loads(quiz_json_string)

            quiz = Quiz.objects.create(
                user_id=request.session["id"],
                module_name=module_name,
                difficulty=difficulty,
                num_questions=num_questions,
                course_id=course_id
            )

            # Save to file
            save_json_to_file(f"./history/{request.session["id"]}/quizzes/", f"{quiz.quiz_id}.json", quiz_data)

            return JsonResponse({"status": "success", "quiz": quiz_data}, status=200)

        except Exception as e:
            print(f"Error generating quiz: {e}")
            return JsonResponse({"status": "error", "message": "An unexpected error occurred."}, status=500)

    

@csrf_exempt
def generate_flashcards_view(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body)
            course_id = data.get("course_id")
            module_name = data.get("module_name")

            if not course_id or not module_name:
                return JsonResponse({"status": "error", "message": "Course ID and module name are required."}, status=400)

            module_file_path = f"./modules/{course_id}/{module_name}"
            extracted_text = extract_text_from_file(module_file_path)
            if not extracted_text:
                return JsonResponse({"status": "error", "message": "Failed to extract text from the module file."}, status=400)

            flashcards = generate_flashcard(extracted_text)
            flashcard_data = json.loads(flashcards)
            if not flashcards:
                return JsonResponse({"status": "error", "message": "Failed to generate flashcards."}, status=500)

            flashcard = Flashcard.objects.create(
                user_id=request.session["id"],
                module_name=module_name,
                course_id=course_id  # Include course_id
            )
            # Save to file
            save_json_to_file(f"./history/{request.session["id"]}/flashcards/", f"{flashcard.flashcard_id}.json", flashcard_data)

            return JsonResponse({"status": "success", "flashcards": flashcards}, status=200)

        except Exception as e:
            print(f"Error generating flashcards: {e}")
            return JsonResponse({"status": "error", "message": "An unexpected error occurred."}, status=500)


@csrf_exempt
def generate_bulletpoints_view(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body)
            course_id = data.get("course_id")
            module_name = data.get("module_name")

            if not course_id or not module_name:
                return JsonResponse({"status": "error", "message": "Course ID and module name are required."}, status=400)

            module_file_path = f"./modules/{course_id}/{module_name}"
            extracted_text = extract_text_from_file(module_file_path)
            if not extracted_text:
                return JsonResponse({"status": "error", "message": "Failed to extract text from the module."}, status=400)

            bullet_points = generate_bulletpoints(extracted_text)
            if not bullet_points:
                return JsonResponse({"status": "error", "message": "Failed to generate bullet points."}, status=500)

            bulletpoint_data=json.loads(bullet_points)

            bulletpoint = Bulletpoint.objects.create(
                user_id=request.session["id"],
                module_name=module_name,
                course_id=course_id  # Include course_id
            )
            # Save to file
            save_json_to_file(f"./history/{request.session["id"]}/bulletpoints/", f"{bulletpoint.bulletpoint_id}.json", bulletpoint_data)

            return JsonResponse({"status": "success", "summary": bullet_points}, status=200)

        except Exception as e:
            print(f"Error generating bullet points: {e}")
            return JsonResponse({"status": "error", "message": "An unexpected error occurred."}, status=500)

    
@csrf_exempt
def generate_fullsummary_view(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body)
            course_id = data.get("course_id")
            module_name = data.get("module_name")

            if not course_id or not module_name:
                return JsonResponse({"status": "error", "message": "Course ID, module name, and summary type are required."}, status=400)

            module_file_path = f"./modules/{course_id}/{module_name}"
            if not os.path.exists(module_file_path):
                return JsonResponse({"status": "error", "message": "Module file does not exist."}, status=404)

            extracted_text = extract_text_from_file(module_file_path)
            if not extracted_text:
                return JsonResponse({"status": "error", "message": "Failed to extract text from the module."}, status=400)

            full_summary = generate_full_summary(extracted_text)
            if not full_summary:
                return JsonResponse({"status": "error", "message": "Failed to generate full summary."}, status=500)

            # Save to file
            fs_data = json.loads(full_summary)

            fullsummary = Fullsummary.objects.create(
                user_id=request.session["id"],
                module_name=module_name,
                course_id=course_id  # Include course_id
            )

            # Save to file
            save_json_to_file(f"./history/{request.session["id"]}/fullsummaries/", f"{fullsummary.fullsummary_id}.json", fs_data)

            return JsonResponse({"status": "success", "summary": full_summary}, status=200)

        except Exception as e:
            print(f"Error generating full summary: {e}")
            return JsonResponse({"status": "error", "message": "An unexpected error occurred."}, status=500)


@csrf_exempt
def generate_steps_for_todos_view(request):
    if request.method == "POST":
        try:
            # Parse the JSON body
            data = json.loads(request.body)
            todo_ids = data.get("todo_ids", [])

            if not todo_ids or not isinstance(todo_ids, list):
                return JsonResponse(
                    {"status": "error", "message": "A list of To-Do IDs is required."},
                    status=400,
                )

            steps_responses = []

            # Fetch and process each To-Do item
            for todo_id in todo_ids:
                try:
                    todo = Todo.objects.get(todo_id=todo_id)
                    file = open(f"./todo/{todo.course.course_id}/{todo_id}.txt", "r")

                    # Generate steps using AI
                    steps = generate_todo_help(todo.assignment_name, file.read())

                    if not steps:
                        steps = ["Failed to generate steps for this To-Do item."]

                    steps_responses.append({
                        "todo_id": todo_id,
                        "assignment_name": todo.assignment_name,
                        "steps": steps
                    })

                except Todo.DoesNotExist:
                    steps_responses.append({
                        "todo_id": todo_id,
                        "error": f"To-Do item with ID {todo_id} not found."
                    })

            # Return the steps for all requested To-Do items
            return JsonResponse(
                {"status": "success", "steps_data": steps_responses},
                status=200,
            )

        except json.JSONDecodeError:
            return JsonResponse(
                {"status": "error", "message": "Invalid JSON format."},
                status=400,
            )
        except Exception as e:
            print(f"Error generating steps for To-Dos: {e}")
            return JsonResponse(
                {"status": "error", "message": "An unexpected error occurred."},
                status=500,
            )
    else:
        return JsonResponse(
            {"status": "error", "message": "Invalid request method."},
            status=405,
        )
    
@csrf_exempt
def get_user_quizzes_view(request):
    if request.method == "GET":
        try:
            # Ensure the user is authenticated
            if not request.user.is_authenticated:
                return JsonResponse(
                    {"status": "error", "message": "User is not authenticated."}, status=401
                )

            # Get the user's ID from the session or request.user
            user_id = request.session.get("id") or request.user.id

            if not user_id:
                return JsonResponse(
                    {"status": "error", "message": "User ID is missing from the session."},
                    status=400
                )

            # Fetch quizzes for the specific user
            quizzes = Quiz.objects.filter(user_id=user_id)

            # Prepare the quizzes data
            quizzes_data = [
                {
                    "quiz_id": quiz.quiz_id,
                    "module_name": quiz.module_name,
                    "difficulty": quiz.difficulty,
                    "num_questions": quiz.num_questions,
                    "course_name": quiz.course.name,
                }
                for quiz in quizzes
            ]

            return JsonResponse({"status": "success", "quizzes": quizzes_data}, status=200)

        except Exception as e:
            print(f"Error fetching user quizzes: {e}")
            return JsonResponse(
                {"status": "error", "message": "Failed to fetch quizzes."},
                status=500
            )
    else:
        return JsonResponse(
            {"status": "error", "message": "Invalid request method."}, status=405
        )
    
@csrf_exempt
def get_quiz_view(request):
    if request.method == "GET":
        try:
            # Extract the quiz ID from the query parameters
            quiz_id = request.GET.get("quizid")
            user_id = request.session.get("id")  # Get the user ID from the session

            if not quiz_id:
                return JsonResponse(
                    {"status": "error", "message": "Quiz ID is required."},
                    status=400
                )
            
            if not user_id:
                return JsonResponse(
                    {"status": "error", "message": "User is not authenticated."},
                    status=401
                )

            # Construct the file path
            file_path = f"./history/{user_id}/quizzes/{quiz_id}.json"

            # Check if the file exists
            if not os.path.exists(file_path):
                return JsonResponse(
                    {"status": "error", "message": "Quiz file not found."},
                    status=404
                )

            # Load the quiz data from the file
            with open(file_path, "r") as file:
                quiz_data = json.load(file)

            return JsonResponse(
                {"status": "success", "quiz": quiz_data},
                status=200
            )

        except Exception as e:
            print(f"Error loading quiz: {e}")
            return JsonResponse(
                {"status": "error", "message": "An unexpected error occurred."},
                status=500
            )
    else:
        return JsonResponse(
            {"status": "error", "message": "Invalid request method."},
            status=405
        )
    
def get_summaries_view(request):
    try:
        # Fetch user ID from the authenticated session
        user_id = request.session["id"]

        # Fetch all summaries for the logged-in user
        flashcards = Flashcard.objects.filter(user_id=user_id)
        bullet_points = Bulletpoint.objects.filter(user_id=user_id)
        full_summaries = Fullsummary.objects.filter(user_id=user_id)

        # Structure the response
        response_data = {
            "status": "success",
            "flashcards": [
                {
                    "flashcard_id": flashcard.flashcard_id,
                    "module_name": flashcard.module_name,
                    "course_name": flashcard.course.name if flashcard.course else None,
                }
                for flashcard in flashcards
            ],
            "bullet_points": [
                {
                    "bulletpoint_id": bulletpoint.bulletpoint_id,
                    "module_name": bulletpoint.module_name,
                    "course_name": bulletpoint.course.name if bulletpoint.course else None,
                }
                for bulletpoint in bullet_points
            ],
            "full_summaries": [
                {
                    "fullsummary_id": fullsummary.fullsummary_id,
                    "module_name": fullsummary.module_name,
                    "course_name": fullsummary.course.name if fullsummary.course else None,
                }
                for fullsummary in full_summaries
            ],
        }

        return JsonResponse(response_data, status=200)

    except Exception as e:
        # Handle any unexpected errors
        return JsonResponse(
            {"status": "error", "message": str(e)}, status=500
        )
    
@csrf_exempt
def get_flashcard_view(request):
    if request.method == "GET":
        try:
            # Extract the flashcard ID from the query parameters
            flashcard_id = request.GET.get("flashcard_id")
            user_id = request.session.get("id")  # Get the user ID from the session

            if not flashcard_id:
                return JsonResponse(
                    {"status": "error", "message": "Flashcard ID is required."},
                    status=400,
                )
            
            if not user_id:
                return JsonResponse(
                    {"status": "error", "message": "User is not authenticated."},
                    status=401,
                )

            # Construct the file path
            file_path = f"./history/{user_id}/flashcards/{flashcard_id}.json"

            # Check if the file exists
            if not os.path.exists(file_path):
                return JsonResponse(
                    {"status": "error", "message": "Flashcard file not found."},
                    status=404,
                )

            # Load the flashcard data from the file
            with open(file_path, "r") as file:
                flashcard_data = json.load(file)

            return JsonResponse(
                {"status": "success", "flashcards": flashcard_data},
                status=200,
            )

        except Exception as e:
            print(f"Error loading flashcard: {e}")
            return JsonResponse(
                {"status": "error", "message": "An unexpected error occurred."},
                status=500,
            )
    else:
        return JsonResponse(
            {"status": "error", "message": "Invalid request method."},
            status=405,
        )
    
@csrf_exempt
def get_bulletpoint_view(request):
    if request.method == "GET":
        try:
            # Extract bullet point ID from query parameters
            bulletpoint_id = request.GET.get("bulletpoint_id")
            user_id = request.session.get("id")  # User ID from the session

            if not bulletpoint_id:
                return JsonResponse(
                    {"status": "error", "message": "Bullet point ID is required."},
                    status=400
                )

            if not user_id:
                return JsonResponse(
                    {"status": "error", "message": "User is not authenticated."},
                    status=401
                )

            # File path for the bullet point summary
            file_path = f"./history/{user_id}/bulletpoints/{bulletpoint_id}.json"

            # Check if the file exists
            if not os.path.exists(file_path):
                return JsonResponse(
                    {"status": "error", "message": "Bullet point file not found."},
                    status=404
                )

            # Read the bullet point data from the file
            with open(file_path, "r") as file:
                bulletpoint_data = json.load(file)

            return JsonResponse(
                {"status": "success", "bullet_points": bulletpoint_data},
                status=200
            )

        except Exception as e:
            print(f"Error loading bullet points: {e}")
            return JsonResponse(
                {"status": "error", "message": "An unexpected error occurred."},
                status=500
            )
    else:
        return JsonResponse(
            {"status": "error", "message": "Invalid request method."},
            status=405
        )
    
@csrf_exempt
def get_fullsummary_view(request):
    if request.method == "GET":
        try:
            # Extract the full summary ID from the query parameters
            fullsummary_id = request.GET.get("fullsummary_id")
            user_id = request.session.get("id")  # Fetch user ID from session

            if not fullsummary_id:
                return JsonResponse(
                    {"status": "error", "message": "Full summary ID is required."},
                    status=400,
                )
            
            if not user_id:
                return JsonResponse(
                    {"status": "error", "message": "User is not authenticated."},
                    status=401,
                )

            # Construct the file path
            file_path = f"./history/{user_id}/fullsummaries/{fullsummary_id}.json"

            # Check if the file exists
            if not os.path.exists(file_path):
                return JsonResponse(
                    {"status": "error", "message": "Full summary file not found."},
                    status=404,
                )

            # Load the summary data from the file
            with open(file_path, "r") as file:
                summary_data = json.load(file)

            return JsonResponse(
                {"status": "success", "summary": summary_data},
                status=200,
            )

        except Exception as e:
            print(f"Error loading full summary: {e}")
            return JsonResponse(
                {"status": "error", "message": "An unexpected error occurred."},
                status=500,
            )
    else:
        return JsonResponse(
            {"status": "error", "message": "Invalid request method."},
            status=405,
        )